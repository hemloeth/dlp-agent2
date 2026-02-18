import logging
import os
import docx
import PyPDF2
import openpyxl
from pptx import Presentation
from dlp_agent.detectors import detect_credit_cards, detect_aadhaar, detect_pan
from dlp_agent.events.sinks import EventSink

class StreamProcessor:
    def __init__(self, config: dict, sinks: list[EventSink] = None):
        self.config = config
        self.sinks = sinks or []
        self.detectors = []
        self.seen_hashes = set() # For deduplication
        self._init_detectors()

    def _init_detectors(self):
        rules = self.config.get('rules', {})
        if rules.get('card', {}).get('enabled', False):
            self.detectors.append(detect_credit_cards)
        if rules.get('aadhaar', {}).get('enabled', False):
            self.detectors.append(detect_aadhaar)
        if rules.get('pan', {}).get('enabled', False):
            self.detectors.append(detect_pan)

    def _get_content_iterator(self, file_path: str):
        """
        Returns an iterator that yields (line_num, content).
        Handles different file types based on extension.
        """
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()

        if ext == '.docx':
            try:
                doc = docx.Document(file_path)
                for i, para in enumerate(doc.paragraphs, 1):
                    yield i, para.text
            except Exception as e:
                logging.error(f"Error reading docx {file_path}: {e}")
        elif ext == '.pdf':
            try:
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for i, page in enumerate(reader.pages, 1):
                        text = page.extract_text()
                        if text:
                            for line_num, line in enumerate(text.splitlines(), 1):
                                # Compounding page and line number for unique tracking
                                yield f"p{i}:l{line_num}", line
            except Exception as e:
                logging.error(f"Error reading pdf {file_path}: {e}")
        elif ext == '.xlsx':
            try:
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                try:
                    for sheet in wb.worksheets:
                        for row_num, row in enumerate(sheet.iter_rows(values_only=True), 1):
                            content = " ".join([str(cell) for cell in row if cell is not None])
                            if content:
                                yield f"{sheet.title}:r{row_num}", content
                finally:
                    wb.close()
            except Exception as e:
                logging.error(f"Error reading xlsx {file_path}: {e}")
        elif ext == '.pptx':
            try:
                prs = Presentation(file_path)
                for i, slide in enumerate(prs.slides, 1):
                    for shape_num, shape in enumerate(slide.shapes, 1):
                        if hasattr(shape, "text"):
                            yield f"s{i}:sh{shape_num}", shape.text
            except Exception as e:
                logging.error(f"Error reading pptx {file_path}: {e}")
        elif ext == '.doc':
            logging.warning(f"File {file_path} is in .doc format. Only .docx is currently supported for Word documents.")
        else:
            # Default text processing
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                for line_num, line in enumerate(f, 1):
                    yield line_num, line

    def process_file(self, file_path: str) -> int:
        """
        Process a file and emit findings to sinks.
        Returns the count of findings.
        """
        findings_count = 0
        try:
            for line_num, line_content in self._get_content_iterator(file_path):
                line_content = line_content.strip()
                if not line_content:
                    continue
                    
                for detector in self.detectors:
                    file_events = detector(line_content)
                    for event in file_events:
                        # Populate source info
                        event.source = {
                            "type": "file",
                            "path": file_path,
                            "line": line_num
                        }
                        
                        # Deduplication check
                        # Key: hash + rule + file + line
                        dedup_key = f"{event.hash}:{event.rule}:{file_path}:{line_num}"
                        
                        if dedup_key in self.seen_hashes:
                            continue
                            
                        self.seen_hashes.add(dedup_key)
                        
                        # Emit to all sinks
                        for sink in self.sinks:
                            sink.emit(event)
                            
                        findings_count += 1
                            
        except Exception as e:
            logging.error(f"Error processing file {file_path}: {str(e)}")
            
        return findings_count
